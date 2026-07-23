from __future__ import annotations

import importlib.util
import json
from pathlib import Path
from types import ModuleType

import pytest


ROOT = Path(__file__).resolve().parents[3]
FLOW_ROOT = ROOT / "flows" / "drm_document_text_extraction_flow"
FLOW_PATH = FLOW_ROOT / "drm_document_text_extraction_flow.json"
COMPONENT_SOURCE = (
    ROOT
    / "components"
    / "drm_document_text_extractor"
    / "drm_document_text_extractor.py"
)


class FakeResponse:
    def __init__(
        self,
        content: bytes,
        *,
        status_code: int = 200,
        content_type: str = "text/plain; charset=utf-8",
        apparent_encoding: str = "utf-8",
    ) -> None:
        self.content = content
        self.status_code = status_code
        self.headers = {
            "Content-Type": content_type,
            "Content-Length": str(len(content)),
        }
        self.apparent_encoding = apparent_encoding
        self.closed = False

    def iter_content(self, chunk_size: int):
        for start in range(0, len(self.content), chunk_size):
            yield self.content[start : start + chunk_size]

    def close(self) -> None:
        self.closed = True


class FakeClient:
    def __init__(self, responses: list[FakeResponse]) -> None:
        self.responses = list(responses)
        self.calls: list[dict] = []

    def post(self, url: str, **kwargs):
        file_tuple = kwargs["files"]["file"]
        self.calls.append(
            {
                "url": url,
                "filename": file_tuple[0],
                "content": file_tuple[1].read(),
                "mime": file_tuple[2],
                "headers": dict(kwargs["headers"]),
                "params": dict(kwargs["params"]),
                "timeout": kwargs["timeout"],
                "verify": kwargs["verify"],
                "allow_redirects": kwargs["allow_redirects"],
                "stream": kwargs["stream"],
            }
        )
        return self.responses.pop(0)


def load_component_module() -> ModuleType:
    spec = importlib.util.spec_from_file_location("drm_document_text_extractor_test", COMPONENT_SOURCE)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_flow() -> dict:
    return json.loads(FLOW_PATH.read_text(encoding="utf-8"))


def decode_handle(value: str) -> dict:
    return json.loads(value.replace("œ", '"'))


def test_package_contract_and_supported_formats() -> None:
    manifest = json.loads((FLOW_ROOT / "manifest.json").read_text(encoding="utf-8"))
    refs = json.loads((FLOW_ROOT / "component_refs.json").read_text(encoding="utf-8"))
    internal = json.loads((FLOW_ROOT / "internal_nodes.json").read_text(encoding="utf-8"))
    component_manifest = json.loads(
        (COMPONENT_SOURCE.parent / "manifest.json").read_text(encoding="utf-8")
    )

    assert manifest["status"] == "user_testing"
    assert manifest["version"] == "0.4.1"
    assert manifest["runtime_ready"] is True
    assert refs == {
        "flow_id": "drm_document_text_extraction_flow",
        "components": [{"id": "drm_document_text_extractor", "version": "0.5.0"}],
    }
    assert internal["nodes"] == []
    file_input = next(item for item in component_manifest["inputs"] if item["name"] == "document_files")
    assert {"pdf", "pptx", "xlsx", "docx", "hwp", "txt", "csv", "png", "jpg", "zip"} <= set(
        file_input["file_types"]
    )
    assert file_input["required"] is False
    assert {item["name"] for item in component_manifest["inputs"]} >= {
        "document_files",
        "file_record",
        "vision_model",
        "drm_api_url",
        "drm_token",
        "employee_no",
        "processing_mode",
    }
    assert {item["name"]: item["types"] for item in component_manifest["outputs"]} == {
        "extracted_text": ["Message"],
        "processed_file": ["Data"],
    }


def test_flow_embeds_component_source_and_connects_message_to_chat_output() -> None:
    flow = load_flow()
    assert flow["name"] == "drm_document_text_extraction_flow"
    assert len(flow["data"]["nodes"]) == 2
    assert len(flow["data"]["edges"]) == 1
    nodes = {node["id"]: node for node in flow["data"]["nodes"]}
    extractor = nodes["DrmDocumentTextExtractor-drmDocument"]
    template = extractor["data"]["node"]["template"]
    assert template["code"]["value"] == COMPONENT_SOURCE.read_text(encoding="utf-8")
    assert template["drm_api_url"]["value"] == ""
    assert template["drm_token"]["value"] == ""
    assert template["employee_no"]["value"] == ""
    assert template["drm_token"]["password"] is True
    assert template["employee_no"]["password"] is True
    assert template["processing_mode"]["value"] == "자동(로컬 우선)"
    assert template["processing_mode"]["real_time_refresh"] is True
    for field_name in ("drm_api_url", "drm_token", "employee_no"):
        assert template[field_name]["required"] is False
    assert template["allow_insecure_http"]["value"] is False
    assert template["verify_tls"]["value"] is True
    assert template["document_files"]["file_path"] in ("", [])
    assert template["document_files"]["required"] is True
    assert template["file_record"]["show"] is False
    assert template["vision_model"]["show"] is False

    edge = flow["data"]["edges"][0]
    assert (edge["source"], edge["target"]) == (
        "DrmDocumentTextExtractor-drmDocument",
        "ChatOutput-drmDocumentText",
    )
    assert decode_handle(edge["sourceHandle"]) == edge["data"]["sourceHandle"]
    assert decode_handle(edge["targetHandle"]) == edge["data"]["targetHandle"]
    assert edge["data"]["sourceHandle"]["name"] == "extracted_text"
    assert edge["data"]["targetHandle"]["fieldName"] == "input_value"


def test_bypass_mode_hides_drm_network_settings() -> None:
    module = load_component_module()
    update_build_config = module.DrmDocumentTextExtractor.update_build_config
    field_names = (
        "drm_api_url",
        "drm_token",
        "employee_no",
        "allow_insecure_http",
        "verify_tls",
        "timeout_seconds",
    )
    config = {name: {"show": True} for name in field_names}
    updated = update_build_config(
        object(), config, module.PROCESSING_MODE_BYPASS_DRM, "processing_mode"
    )
    assert all(updated[name]["show"] is False for name in field_names)

    updated = update_build_config(
        object(), updated, module.PROCESSING_MODE_AUTO, "processing_mode"
    )
    assert all(updated[name]["show"] is True for name in field_names)


def test_extracts_multiple_files_with_original_request_contract(tmp_path: Path) -> None:
    module = load_component_module()
    pdf = tmp_path / "guide.pdf"
    xlsx = tmp_path / "cost.xlsx"
    pdf.write_bytes(b"%PDF fake")
    xlsx.write_bytes(b"PK fake-xlsx")
    first_text = "PDF 본문입니다."
    second_text = "엑셀 본문입니다."
    responses = [
        FakeResponse(first_text.encode("utf-8")),
        FakeResponse(
            second_text.encode("cp949"),
            content_type="text/plain; charset=cp949",
            apparent_encoding="cp949",
        ),
    ]
    client = FakeClient(responses)

    result = module.extract_document_texts(
        [str(pdf), str(xlsx)],
        "https://drm.example.internal/DRM/decrypt/text",
        "TOKEN_VALUE",
        "X000000",
        transport=client,
    )

    assert result["success"] is True
    assert [item["file_name"] for item in result["data"]] == ["guide.pdf", "cost.xlsx"]
    assert [item["text"] for item in result["data"]] == [first_text, second_text]
    assert result["meta"]["file_count"] == 2
    assert result["meta"]["local_file_count"] == 0
    assert result["meta"]["drm_file_count"] == 2
    assert result["meta"]["network_called"] is True
    assert all(response.closed for response in responses)
    assert [call["content"] for call in client.calls] == [b"%PDF fake", b"PK fake-xlsx"]
    for call in client.calls:
        assert call["mime"] == "application/octet-stream"
        assert call["headers"] == {"Authorization": "Bearer TOKEN_VALUE"}
        assert call["params"] == {"empNo": "X000000"}
        assert call["timeout"] == 180
        assert call["verify"] is True
        assert call["allow_redirects"] is False
        assert call["stream"] is True

    message = module.format_extraction_message(result)
    assert "[FILE 1/2] guide.pdf" in message
    assert first_text in message
    assert "[FILE 2/2] cost.xlsx" in message
    assert second_text in message


def test_auto_mode_extracts_plain_office_and_text_files_without_drm(tmp_path: Path) -> None:
    module = load_component_module()

    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    docx_path = tmp_path / "plain.docx"
    document = Document()
    document.add_paragraph("일반 Word 본문")
    document.save(docx_path)

    pptx_path = tmp_path / "plain.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "일반 PowerPoint 제목"
    slide.placeholders[1].text = "발표 본문"
    presentation.save(pptx_path)

    xlsx_path = tmp_path / "plain.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "계획"
    worksheet.append(["항목", "금액"])
    worksheet.append(["설비", 1200])
    workbook.save(xlsx_path)
    workbook.close()

    txt_path = tmp_path / "plain.txt"
    txt_path.write_text("일반 텍스트 본문", encoding="utf-8")

    client = FakeClient([])
    result = module.extract_document_texts(
        [docx_path, pptx_path, xlsx_path, txt_path],
        processing_mode=module.PROCESSING_MODE_AUTO,
        transport=client,
    )

    assert client.calls == []
    assert result["meta"]["local_file_count"] == 4
    assert result["meta"]["drm_file_count"] == 0
    assert result["meta"]["network_called"] is False
    assert [item["processing_path"] for item in result["data"]] == ["local"] * 4
    combined = "\n".join(item["text"] for item in result["data"])
    assert "일반 Word 본문" in combined
    assert "일반 PowerPoint 제목" in combined
    assert "설비\t1200" in combined
    assert "일반 텍스트 본문" in combined


def test_bypass_mode_never_calls_drm_and_rejects_nonlocal_direct_format(tmp_path: Path) -> None:
    module = load_component_module()
    text_file = tmp_path / "plain.csv"
    text_file.write_text("name,value\nalpha,1", encoding="utf-8")
    client = FakeClient([])

    result = module.extract_document_texts(
        text_file,
        processing_mode=module.PROCESSING_MODE_BYPASS_DRM,
        transport=client,
    )
    assert result["data"][0]["processing_path"] == "local"
    assert result["meta"]["network_called"] is False
    assert client.calls == []

    legacy = tmp_path / "legacy.xls"
    legacy.write_bytes(b"legacy workbook")
    with pytest.raises(ValueError, match="DRM 미사용 모드에서 로컬 텍스트 추출에 실패"):
        module.extract_document_texts(
            legacy,
            processing_mode=module.PROCESSING_MODE_BYPASS_DRM,
            transport=client,
        )
    assert client.calls == []


def test_http_policy_fails_before_upload(tmp_path: Path) -> None:
    module = load_component_module()
    document = tmp_path / "review.docx"
    document.write_bytes(b"PK fake-docx")
    client = FakeClient([FakeResponse(b"unused")])

    with pytest.raises(ValueError, match="HTTP로 보낼 수 없습니다"):
        module.extract_document_texts(
            str(document),
            "http://drm.example.internal/DRM/decrypt/text",
            "TOKEN_VALUE",
            "X000000",
            transport=client,
        )
    assert client.calls == []


def test_http_can_be_explicitly_enabled_for_closed_network(tmp_path: Path) -> None:
    module = load_component_module()
    document = tmp_path / "legacy.xls"
    document.write_bytes(b"legacy-xls")
    response = FakeResponse("셀 내용".encode("utf-8"))
    client = FakeClient([response])

    result = module.extract_document_texts(
        str(document),
        "http://drm.example.internal/DRM/decrypt/text",
        "TOKEN_VALUE",
        "X000000",
        allow_insecure_http=True,
        transport=client,
    )
    assert result["data"][0]["text"] == "셀 내용"
    assert client.calls[0]["verify"] is False


def test_unsupported_format_is_skipped_and_size_fails_before_upload(tmp_path: Path) -> None:
    module = load_component_module()
    text_file = tmp_path / "archive.zip"
    text_file.write_bytes(b"not supported")
    client = FakeClient([FakeResponse(b"unused")])
    skipped = module.extract_document_texts(
        str(text_file),
        "https://drm.example.internal/DRM/decrypt/text",
        "TOKEN_VALUE",
        "X000000",
        transport=client,
    )
    assert skipped["data"][0]["processing_path"] == "skipped_unsupported"
    assert skipped["data"][0]["drm_status"] == "not_applicable"
    assert "지원하지 않는 첨부파일 형식" in skipped["data"][0]["text"]
    assert skipped["meta"]["skipped_file_count"] == 1
    assert skipped["meta"]["network_called"] is False

    pdf = tmp_path / "large.pdf"
    pdf.write_bytes(b"x" * (1024 * 1024 + 1))
    with pytest.raises(ValueError, match="개별 크기 제한"):
        module.extract_document_texts(
            str(pdf),
            "https://drm.example.internal/DRM/decrypt/text",
            "TOKEN_VALUE",
            "X000000",
            max_file_size_mb=1,
            transport=client,
        )
    assert client.calls == []


def test_http_error_does_not_expose_response_body_or_secret(tmp_path: Path) -> None:
    module = load_component_module()
    document = tmp_path / "review.pptx"
    document.write_bytes(b"PK fake-pptx")
    secret = "TOP_SECRET_TOKEN"
    response = FakeResponse(
        b"sensitive server details TOP_SECRET_TOKEN",
        status_code=500,
    )
    client = FakeClient([response])

    with pytest.raises(module.DrmTextExtractionError) as caught:
        module.extract_document_texts(
            str(document),
            "https://drm.example.internal/DRM/decrypt/text",
            secret,
            "X000000",
            transport=client,
        )
    message = str(caught.value)
    assert "HTTP 500" in message
    assert secret not in message
    assert "sensitive server details" not in message
    assert response.closed is True


def test_samples_and_exports_do_not_contain_real_credentials() -> None:
    texts = []
    for path in [
        FLOW_ROOT / "README.md",
        FLOW_ROOT / "samples" / "INPUT_EXAMPLE.md",
        COMPONENT_SOURCE.parent / "USAGE_GUIDE.md",
        FLOW_PATH,
    ]:
        texts.append(path.read_text(encoding="utf-8"))
    combined = "\n".join(texts)
    assert "TOKEN_INFO" not in combined
    assert "X123456" not in combined
    assert "drmtest.com" not in combined
