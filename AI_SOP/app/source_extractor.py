from __future__ import annotations

import csv
import io
from pathlib import Path


def extract_text(filename: str, media_type: str, data: bytes, max_chars: int = 120_000) -> str:
    suffix = Path(filename).suffix.lower()
    text = ""
    try:
        if suffix in {".txt", ".md", ".csv"}:
            text = data.decode("utf-8-sig", errors="replace")
        elif suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
        elif suffix == ".docx":
            from docx import Document

            document = Document(io.BytesIO(data))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        elif suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            rows: list[str] = []
            for worksheet in workbook.worksheets:
                rows.append(f"[Sheet: {worksheet.title}]")
                for row in worksheet.iter_rows(values_only=True):
                    rows.append("\t".join("" if value is None else str(value) for value in row))
            text = "\n".join(rows)
        elif suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(io.BytesIO(data))
            slides: list[str] = []
            for index, slide in enumerate(presentation.slides, start=1):
                slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                slides.append(f"[Slide {index}]\n" + "\n".join(slide_text))
            text = "\n\n".join(slides)
    except Exception:
        return ""
    return text[:max_chars]

